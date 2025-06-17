# lambda/index.py
import base64
import boto3
import io
import json
import os
import pptx
import re

# Lambda コンテキストからリージョンを抽出する関数
def extract_region_from_arn(arn):
    # ARN 形式: arn:aws:lambda:region:account-id:function:function-name
    match = re.search('arn:aws:lambda:([^:]+):', arn)
    if match:
        return match.group(1)
    return "us-east-1"  # デフォルト値

# グローバル変数としてクライアントを初期化（初期値）
bedrock_client = None

# モデルID
MODEL_ID = os.environ.get("MODEL_ID", "us.amazon.nova-lite-v1:0")

# .pptxからテキスト抽出
def extract_text_from_pptx(pptx_bytes):
    prs = pptx.Presentation(io.BytesIO(pptx_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def lambda_handler(event, context):
    try:
        # コンテキストから実行リージョンを取得し、クライアントを初期化
        global bedrock_client
        if bedrock_client is None:
            region = extract_region_from_arn(context.invoked_function_arn)
            bedrock_client = boto3.client('bedrock-runtime', region_name=region)
            print(f"Initialized Bedrock client in region: {region}")
        
        print("Received event:", json.dumps(event))
        
        # Cognitoで認証されたユーザー情報を取得
        user_info = None
        if 'requestContext' in event and 'authorizer' in event['requestContext']:
            user_info = event['requestContext']['authorizer']['claims']
            print(f"Authenticated user: {user_info.get('email') or user_info.get('cognito:username')}")
        
        # リクエストボディの解析
        body = json.loads(event['body'])
        # 入力ファイルのデコードとテキスト抽出
        pptx_base64 = body['pptxFile']
        txt_base64 = body['txtFile']

        if not pptx_base64 or not txt_base64:
            raise Exception("Both pptxFile and txtFile must be provided in base64 format.")
        
        pptx_bytes = base64.b64decode(pptx_base64)
        txt_bytes = base64.b64decode(txt_base64)
        original_txt = txt_bytes.decode('utf-8')

        # pptxからの正しい内容の抽出
        reference_text = extract_text_from_pptx(pptx_bytes)
        
        # 修正指示を構成してLLMに送信
        instruction_message = f"""以下の2つのファイルをもとに、テキストの校正と自然な文章への修正を行ってください。
        1. 講義資料としてのパワーポイント
        2. 音声から生成された書き起こしの修正対象のテキスト
        【依頼内容】  
        - 音声テキストは講義資料に沿って作成されているため、両方を参照しながら、読みやすく自然な日本語文章に修正してください。  
        - 不自然な「どもり」や「言い直し」などの話し言葉は適切に取り除き、滑らかな文章にしてください。  
        - パワーポイントの資料に含まれる内容が不足している場合は適宜補足し、整合性を持たせてください。  
        - もし、音声テキストだけで十分に修正が可能な場合は、資料の参照は必須ではありませんが、基本的には両方を活用してください。  
        【最終成果物】  
        自然でわかりやすい文章として校正されたテキストを出力してください。英語の場合は英語で、その他の言語の場合はその言語で出力してください。
        修正後のテキストのみ出力してください。  
        【入力内容】  
        --- 講義資料としてのパワーポイント ---
        {reference_text}
        --- 音声から生成された書き起こしの修正対象のテキスト ---
        {original_txt}  
        """
        
        # invoke_model用のリクエストペイロード
        request_payload = {
            "messages": [{
                "role": "user",
                "content": [{"text": instruction_message}]
            }],
            "inferenceConfig": {
                "maxTokens": 512,
                "stopSequences": [],
                "temperature": 0.7,
                "topP": 0.9
            }
        }
        
        print("Calling Bedrock invoke_model API with payload:", json.dumps(request_payload))
        
        # invoke_model APIを呼び出し
        response = bedrock_client.invoke_model(
            modelId=MODEL_ID,
            body=json.dumps(request_payload),
            contentType="application/json"
        )
        
        # レスポンスを解析
        response_body = json.loads(response['body'].read())
        print("Bedrock response:", json.dumps(response_body, default=str))
        
        # 応答の検証
        if not response_body.get('output') or not response_body['output'].get('message') or not response_body['output']['message'].get('content'):
            raise Exception("No response content from the model")
        
        # アシスタントの応答を取得
        assistant_response = response_body['output']['message']['content'][0]['text']
        
        # 成功レスポンスの返却
        return {
            "statusCode": 200,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                "Access-Control-Allow-Methods": "OPTIONS,POST"
            },
            "body": json.dumps({
                "success": True,
                "response": assistant_response
            })
        }
        
    except Exception as error:
        print("Error:", str(error))
        
        return {
            "statusCode": 500,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token",
                "Access-Control-Allow-Methods": "OPTIONS,POST"
            },
            "body": json.dumps({
                "success": False,
                "error": str(error)
            })
        }
    